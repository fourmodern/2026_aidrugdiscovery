#!/usr/bin/env python3
"""학회 구두발표용 슬라이드 (16:9, 12분 · 12장).

`make_docs.py` 가 만드는 pptx 는 보고서를 기계적으로 옮긴 것이라 발표에 쓸 수 없다.
발표는 문서 낭독이 아니라 **하나의 논증**이다. 그래서 별도로 만든다.

원칙
  - 한 장에 메시지 하나. 슬라이드 제목이 그 메시지 자체다 (주제어가 아니라 문장).
  - 숫자는 크게. 뒷줄에서 읽혀야 한다.
  - 그림은 꽉 채운다. 캡션은 슬라이드 제목이 대신한다.
  - 발표자 노트에 말할 내용을 적는다 — 슬라이드에 적지 않는다.
"""
from __future__ import annotations
import argparse, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
SR = ROOT / "sample_run"
FIG = SR / "report" / "figures_controlled"
FONT = "Pretendard"
W, H = 338.7, 190.5                    # 16:9 (13.333 × 7.5 in)

INK    = RGBColor(0x1A, 0x1A, 0x1A)
NAVY   = RGBColor(0x0B, 0x2A, 0x4A)
BLUE   = RGBColor(0x00, 0x72, 0xB2)
ORANGE = RGBColor(0xD5, 0x5E, 0x00)
GREEN  = RGBColor(0x00, 0x9E, 0x73)
GREY   = RGBColor(0x5A, 0x63, 0x6B)
LGREY  = RGBColor(0xF3, 0xF5, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PALE   = RGBColor(0x9F, 0xC5, 0xE0)


def load(n):
    p = SR / n
    if not p.exists(): return None
    d = json.loads(p.read_text()); return d.get("result", d)


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = Mm(W), Mm(H)
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def box(self, s, x, y, w, h, fill=None, line=None, lw=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE):
        sh = s.shapes.add_shape(shape, Mm(x), Mm(y), Mm(w), Mm(h))
        if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
        else: sh.fill.background()
        if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
        else: sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def text(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space=0):
        tb = s.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Mm(1); tf.margin_top = tf.margin_bottom = 0
        paras = runs if isinstance(runs[0], list) else [runs]
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(space)
            for t, sz, colr, *rest in para:
                r = p.add_run(); r.text = t
                r.font.size = Pt(sz); r.font.color.rgb = colr
                r.font.bold = bool(rest and rest[0]); r.font.name = FONT
        return tb

    def slide(self, title, kicker="", note=""):
        """제목은 주제어가 아니라 그 장에서 하려는 말 자체다."""
        s = self.prs.slides.add_slide(self.blank); self.n += 1
        self.box(s, 0, 0, W, 2.2, fill=BLUE, shape=MSO_SHAPE.RECTANGLE)
        if kicker:
            self.text(s, 14, 6, 120, 6, [(kicker, 12, BLUE, True)])
        self.text(s, 14, 12, W - 28, 16, [(title, 25, NAVY, True)])
        self.text(s, W - 26, H - 11, 16, 6, [(str(self.n), 12, GREY)],
                  align=PP_ALIGN.RIGHT)
        if note:
            s.notes_slide.notes_text_frame.text = note
        return s

    def pic(self, s, name, x, y, w, maxh=None):
        f = FIG / name
        if not f.exists():
            self.box(s, x, y, w, maxh or 60, fill=LGREY, line=GREY)
            return
        p = s.shapes.add_picture(str(f), Mm(x), Mm(y), width=Mm(w))
        if maxh and p.height / 36000 > maxh:
            r = maxh / (p.height / 36000)
            p.height = int(p.height * r); p.width = int(p.width * r)
            p.left = Mm(x + (w - p.width / 36000) / 2)
        return p

    def bignum(self, s, x, y, w, val, lab, colr=ORANGE, sz=54):
        self.text(s, x, y, w, 22, [(val, sz, colr, True)], align=PP_ALIGN.CENTER)
        self.text(s, x, y + 20, w, 10, [(lab, 13, GREY)], align=PP_ALIGN.CENTER)

    def save(self, out):
        self.prs.save(str(out)); return self.n


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(SR / "report" / "docs" / "talk_16x9.pptx"))
    a = ap.parse_args()
    ds, dk, an = load("dataset_controlled.json"), load("docking_controlled.json"), \
                 load("analysis_controlled.json")
    tc, sw = load("terms_controlled.json"), load("exhaustiveness_sweep.json")
    prot, col = load("protonation_test.json"), load("collapse_diagnosis.json")
    oldp = load("old_set_new_protocol.json")
    if not (ds and dk and an):
        raise SystemExit("산출 파일이 없다 — 슬라이드를 만들지 않는다.")
    T = an["arms"]["top_pose"]; PV = T["pose_validity"]; PVI = PV["interpretable_only"]
    wb = T["within_similarity_bin"]
    ACI = T["auc_ci95"]; NSEN = T["near_sensitivity"]; ctrl = dk["control_redock"]
    d = Deck()

    # ── 1 표지 ──────────────────────────────────────────────────────
    s = d.prs.slides.add_slide(d.blank); d.n += 1
    d.box(s, 0, 0, W, H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    d.text(s, 24, 52, W - 48, 30,
           [("골격 유사성과 역가를 직교화한\nPDE5A 도킹 재평가", 40, WHITE, True)],
           align=PP_ALIGN.CENTER)
    d.text(s, 24, 100, W - 48, 12,
           [("선행 상관의 재현 실패, 그리고 자세 타당도가 정하는 해석의 상한", 18, PALE)],
           align=PP_ALIGN.CENTER)
    d.text(s, 24, 132, W - 48, 20,
           [[(f"smina / AutoDock Vina 1.1.2  ·  PDB 1UDT  ·  ChEMBL PDE5A  n = {an['n']}",
              13, PALE)],
            [("전 산출물·코드·리뷰 기록 공개 — github.com/fourmodern/2026_aidrugdiscovery",
              13, WHITE, True)]], align=PP_ALIGN.CENTER, space=4)
    s.notes_slide.notes_text_frame.text = (
        "12분. 배분: 2번 45초 · 3번 45초 · 4번 70초 · 5번 60초 · 6번 70초 · 7·8번 각 60초 · "
        "9번 70초 · 10번 70초 · 11번 75초 · 12번 60초 · 13번 35초.\n"
        "구성: 문제 → 사전 합격선 → 답 → 설계 → (자세 갈래 6-8) → (상관 갈래 9-11) → 정리.")

    # ── 2 문제 ──────────────────────────────────────────────────────
    s = d.slide("홀로 구조를 쓰면 점수와 골격 닮음이 같은 방향을 가리킨다", "문제",
                "표적 좌표는 대개 어떤 리간드가 결합한 채로 풀린 holo 구조입니다. 그 리간드를 "
                "닮은 화합물은 두 경로로 유리해집니다. 포켓이 그 리간드에 맞춰 유도적합돼 "
                "있으니 기하학적으로 잘 맞고, 같은 계열을 최적화한 산물이니 실제로도 강력합니다. "
                "방향이 같으니 상관은 나옵니다. 그런데 그게 채점 능력인지 골격 닮음인지 "
                "구분이 안 됩니다.\n"
                "[다음으로] 그래서 우리는 먼저 무엇을 성공으로 볼지부터 정했습니다.")
    d.pic(s, "fig00_problem.png", 52, 32, W - 104, maxh=142)

    # ── 3 사전 합격선 (신설) ────────────────────────────────────────
    s = d.slide("사전에 정한 합격선은 넷이었다", "약속",
                "숫자를 보여드리기 전에 합격선을 먼저 말씀드립니다. 뒤에 나오는 판정은 전부 "
                "이 표 대비입니다. 목적 4는 원래 계획에 없었는데, 재도킹 대조가 실패하면서 "
                "추가됐습니다.\n"
                "[다음으로] 결과부터 말씀드리면, 넷 중 하나도 넘지 못했습니다.")
    rows = [("1", "역가를 정량 예측하는가", "교차검증 Q²", "Q² ≥ 0.3"),
            ("2", "강한 화합물을 선별하는가", "순위상관 · ROC-AUC", "p < 0.05  &  AUC ≥ 0.7"),
            ("3", "골격 통제 후에도 남는가", "Tanimoto 편상관", "원상관의 60% 이상 유지"),
            ("4", "대조 실패 시 무엇이 병목인가", "C1 / C2 분리 + 통제 실험", "후보를 실험으로 제거")]
    for j, (lab, x, w) in enumerate((("목적", 40, 118), ("판정 지표", 162, 82),
                                     ("사전 합격선", 250, 74))):
        d.text(s, x, 40, w, 8, [(lab, 14, GREY, True)],
               align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
    for i, (n_, obj, met, cri) in enumerate(rows):
        yy = 50 + i * 23
        d.box(s, 20, yy, W - 40, 21, fill=(LGREY if i % 2 == 0 else WHITE),
              shape=MSO_SHAPE.RECTANGLE)
        d.box(s, 24, yy + 4, 12, 13, fill=BLUE, shape=MSO_SHAPE.OVAL)
        d.text(s, 24, yy + 4, 12, 13, [(n_, 15, WHITE, True)],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        d.text(s, 40, yy + 4, 118, 13, [(obj, 17, INK)], anchor=MSO_ANCHOR.MIDDLE)
        d.text(s, 162, yy + 4, 82, 13, [(met, 15, GREY)],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        d.text(s, 250, yy + 4, 74, 13, [(cri, 16, BLUE, True)],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    d.text(s, 20, 146, W - 40, 10,
           [("사후성 고지 — 이 기준은 문헌 관행에서 가져왔고 데이터를 본 뒤 명문화했다 (논문 §2.6)",
             13, GREY)], align=PP_ALIGN.CENTER)

    # ── 4 결론부터 ──────────────────────────────────────────────────
    s = d.slide("결론부터 — 넷 중 하나도 넘지 못했다", "핵심",
                "전체 순위상관이 -0.079, p는 0.32입니다. Q²는 -0.018로 기준 0.3에, AUC는 "
                "0.574로 기준 0.7에 미달합니다. 중요한 건 아래 문장입니다 — 이건 귀무를 "
                "채택한 게 아닙니다. 신뢰구간 상한이 +0.08이니 약한 관계까지 배제하지는 "
                "못합니다.\n"
                "[다음으로] 이 숫자가 나온 설계를 보시겠습니다.")
    d.box(s, 20, 40, W - 40, 34, fill=RGBColor(0xFD, 0xF0, 0xE6), line=ORANGE, lw=2)
    d.text(s, 26, 47, W - 52, 14,
           [(f"ρ = {T['spearman']:+.3f}    (95% CI {T['ci95'][0]:+.2f} ~ "
             f"{T['ci95'][1]:+.2f},   p = {T['perm_p']})", 34, ORANGE, True)],
           align=PP_ALIGN.CENTER)
    for i, (v, lab, crit) in enumerate([
            (f"{tc['models']['single_vina_score']['Q2_loo']:+.3f}" if tc else "—",
             "교차검증 Q²", "기준 0.3  →  미달"),
            (f"{ACI['all']['auc']}", "ROC-AUC 전체", "기준 0.7  →  미달"),
            (f"{T['perm_p']}", "순열 검정 p", "기준 0.05  →  미달")]):
        x = 20 + i * ((W - 40) / 3)
        d.text(s, x, 84, (W - 40) / 3, 20, [(v, 48, ORANGE, True)], align=PP_ALIGN.CENTER)
        d.text(s, x, 106, (W - 40) / 3, 8, [(lab, 15, INK)], align=PP_ALIGN.CENTER)
        d.text(s, x, 115, (W - 40) / 3, 8, [(crit, 13, GREY)], align=PP_ALIGN.CENTER)
    d.box(s, 20, 130, W - 40, 30, fill=LGREY, line=NAVY, lw=1.4)
    d.text(s, 26, 135, W - 52, 22,
           [[("“신호 없음” 은 귀무를 채택한 것이 아니다", 19, NAVY, True)],
            [(f"신뢰구간 상한이 {T['ci95'][1]:+.2f} 이므로 |ρ| ≲ {abs(T['ci95'][0]):.2f} 인 "
              f"관계는 이 데이터로 배제하지 못한다. 기각하지 못한 것이지 없다고 보인 것이 아니다.",
              15, INK)]], space=3)

    # ── 5 설계 ──────────────────────────────────────────────────────
    s = d.slide("역가와 골격 유사도를 9칸 격자로 직교화했다", "설계",
                "ChEMBL 전수를 받아 화합물당 중앙값으로 집계하고, 역가 3구간과 실데나필 "
                "Tanimoto 3구간의 격자에서 균등 추출했습니다. 교란 상관이 0.281에서 0.175로 "
                "떨어집니다. 완전한 직교화는 화학적으로 불가능합니다 — 실데나필과 매우 닮았는데 "
                "약한 화합물이 드물어서 오른쪽 아래 두 칸이 얇습니다.\n"
                "[다음으로] 그런데 이 결과를 해석하기 전에 확인할 게 있었습니다. 재도킹 대조입니다.")
    d.pic(s, "fig02_dataset.png", 12, 38, W - 24, maxh=92)
    d.text(s, 20, 136, W - 40, 18,
           [[(f"ChEMBL PDE5A 활성 {ds['records_scanned']} 레코드 전수 → 화합물당 중앙값 "
              f"({ds['pool_size']}종) → 9칸 격자 균등 추출 n = {an['n']}", 15, INK)],
            [(f"교란 상관 Pearson {ds['confound_pearson_r_potency_vs_similarity']:+.3f}"
              f"   (역가로만 층화한 옛 설계 +0.281)", 17, GREEN, True)]],
           align=PP_ALIGN.CENTER, space=3)

    # ── 6 대조 실패 + 81% ───────────────────────────────────────────
    s = d.slide(f"재도킹 대조를 쪼개니 채점이 실패했다 — {PVI['n']}개 중 "
                f"{1 - PVI['frac_under_threshold']:.0%}가 같은 실패다", "해석의 상한",
                "C1은 상위 모드 중 최선, C2는 점수 1위 자세입니다. 합쳐 하나로 보고하면 탐색 "
                "실패인지 채점 실패인지 알 수 없습니다. 공결정 리간드에서 C2가 실패했는데, "
                "이건 한 건짜리 일화가 아닙니다. 163개 전체로 재면 중앙값 4.50 Å, 2 Å 기준을 "
                "넘는 게 19%뿐입니다. 즉 앞으로 보실 상관은 대부분 틀린 자세에 매긴 점수에 "
                "대한 것입니다. 이게 이 연구에서 가장 큰 제약입니다.\n"
                "[다음으로] 그럼 왜 채점이 실패했나. 후보를 하나씩 지워 보겠습니다.")
    d.pic(s, "fig03_control.png", 10, 36, 176, maxh=98)
    for i, (nm, v, ok_, c) in enumerate([
            ("C1  샘플링", f"{ctrl['rmsd_best_of_modes_angstrom']:.2f} Å", "PASS", GREEN),
            ("C2  채점", f"{ctrl['rmsd_top_pose_angstrom']:.2f} Å", "FAIL", ORANGE)]):
        d.box(s, 196, 40 + i * 22, 132, 19, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        d.text(s, 200, 43 + i * 22, 124, 13, [(f"{nm}   {v}   {ok_}", 18, WHITE, True)],
               align=PP_ALIGN.CENTER)
    d.box(s, 196, 88, 132, 46, fill=RGBColor(0xFD, 0xF0, 0xE6), line=ORANGE, lw=1.6)
    d.text(s, 200, 92, 124, 40,
           [[("그리고 이건 n=1 이 아니다", 16, ORANGE, True)],
            [(f"해석 가능한 대역 {PVI['n']}건의 MCS-RMSD 중앙값 "
              f"{PVI['median_rmsd']:.2f} Å, 2 Å 기준 통과 "
              f"{PVI['frac_under_threshold']:.1%}", 14, INK)],
            [(f"→ 뒤에 나오는 상관은 대부분 틀린 자세에 매긴 점수에 대한 것이다", 14, ORANGE, True)]],
           space=3)
    d.text(s, 10, 140, W - 20, 14,
           [("공결정 리간드 재도킹의 C2 실패는 예외가 아니라 전체의 축소판이다 — "
             "같은 기준, 같은 실패, 163배 규모", 15, NAVY, True)], align=PP_ALIGN.CENTER)

    # ── 7 탐색 깊이 ─────────────────────────────────────────────────
    s = d.slide("탐색을 32배 늘려도 채점은 미동도 없었다", "C2 실패 원인 배제 ①",
                "깊이 4에서 128까지, 시드 3개. C1은 통과율 0에서 100%로 완전히 고쳐지는데 "
                "C2는 전 구간 0%이고 최고 점수는 아예 변하지 않습니다. 해석하면 — 탐색은 가장 "
                "얕은 깊이에서 이미 전역 최소를 찾았고, 그 최소가 결정 자세에서 8.4 Å 떨어진 "
                "자리에 있습니다. 탐색으로는 고칠 수 없다는 뜻입니다.\n"
                "[다음으로] 두 번째 후보는 프로토네이션이었습니다.")
    d.pic(s, "fig08_exhaustiveness.png", 12, 38, W - 24, maxh=104)
    d.text(s, 20, 148, W - 40, 12,
           [[("C1 통과율 0% → 100%", 17, GREEN, True), ("        ", 17, INK),
             (f"C2 전 구간 0%, 변동 {sw['diagnosis']['c2_range_angstrom']} Å", 17, ORANGE, True),
             ("        ", 17, INK),
             (f"최고 점수 변동 {sw['diagnosis']['score_range_kcal']} kcal/mol", 17, ORANGE, True)]],
           align=PP_ALIGN.CENTER)

    # ── 8 프로토네이션 ──────────────────────────────────────────────
    s = d.slide("프로토네이션도 원인이 아니었다 — 그리고 그럴 수 없었다", "C2 실패 원인 배제 ②",
                "수용체는 pH 7.4로 프로토네이션했는데 리간드는 중성 그대로였습니다. 실데나필의 "
                "피페라진은 생리적 pH에서 양성자화되니 이 비대칭이 원인일 수 있었습니다. 2×2로 "
                "맞춰 봤더니 네 조건 점수가 전부 -9.60으로 같습니다. 이유까지 확인했습니다 — "
                "기본 Vina 함수에 정전기 항이 없습니다. 타이핑은 실제로 바뀌었는데 그 질소가 "
                "결합 거리 밖이라 수소결합 항이 안 움직였습니다.\n"
                "[다음으로] 자세 갈래는 여기까지입니다. 이제 본 결과인 상관으로 갑니다.")
    rows2 = [("수용체 pH7.4 · 리간드 pH7.4  (옳은 짝)", "recpH74_ligpH74"),
             ("수용체 pH7.4 · 리간드 중성  (현행)", "recpH74_ligneutral"),
             ("수용체 중성 · 리간드 중성", "recneutral_ligneutral"),
             ("수용체 중성 · 리간드 pH7.4", "recneutral_ligpH74")]
    by = {x["condition"]: x for x in (prot or {}).get("summary", [])}
    d.text(s, 22, 38, 150, 8, [("조건", 14, GREY, True)])
    for j, lab in enumerate(("C1 (Å)", "C2 (Å)", "점수", "C2 통과")):
        d.text(s, 180 + j * 38, 38, 36, 8, [(lab, 14, GREY, True)], align=PP_ALIGN.CENTER)
    for i, (lab, key) in enumerate(rows2):
        r = by.get(key, {}); yy = 48 + i * 15
        d.box(s, 20, yy, W - 40, 13.5, fill=(LGREY if i % 2 == 0 else WHITE),
              shape=MSO_SHAPE.RECTANGLE)
        d.text(s, 24, yy + 1.5, 154, 11, [(lab, 15, INK, i == 0)], anchor=MSO_ANCHOR.MIDDLE)
        for j, v in enumerate((f"{r.get('c1_best_rmsd_mean', 0):.2f}",
                               f"{r.get('c2_top_rmsd_mean', 0):.2f}",
                               f"{r.get('top_score_mean', 0):.2f}",
                               f"{r.get('c2_pass_rate', 0):.0%}")):
            d.text(s, 180 + j * 38, yy + 1.5, 36, 11,
                   [(v, 16, ORANGE if j == 3 else INK, True)],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    d.box(s, 20, 114, W - 40, 30, fill=RGBColor(0xFD, 0xF0, 0xE6), line=ORANGE, lw=1.4)
    d.text(s, 26, 118, W - 52, 24,
           [[("기본 Vina 채점 함수에는 정전기 항이 없다.", 18, ORANGE, True)],
            [("형식전하는 점수에 닿지 못하고, 바뀐 수소결합 타이핑도 그 질소가 결합 거리 "
              "밖이라 항 값이 움직이지 않았다. 원리적으로 원인이 될 수 없다.", 15, INK)]], space=3)
    d.text(s, 20, 150, W - 40, 10,
           [("여기까지가 자세 갈래 — C2 실패의 두 후보를 지웠다. 이제 상관 갈래로 넘어간다.",
             14, NAVY, True)], align=PP_ALIGN.CENTER)

    # ── 9 주 결과 ───────────────────────────────────────────────────
    s = d.slide("교란을 통제해도, 대역을 나눠도 신호는 없다", "결과 ①",
                "골격과 크기를 각각, 그리고 동시에 통제했습니다. 원상관 -0.079가 통제 후에도 "
                "-0.088에서 -0.105 사이로 거의 그대로입니다. 유사도 대역으로 쪼개도 far와 mid는 "
                "0 근처입니다. 막대는 95% 신뢰구간이고 전부 0을 포함합니다.\n"
                "[다음으로] 딱 한 대역만 다릅니다.")
    d.pic(s, "fig05_forest.png", 22, 32, W - 44, maxh=126)
    d.text(s, 20, 162, W - 40, 12,
           [[("전체 ", 16, GREY), (f"{T['spearman']:+.3f}", 19, ORANGE, True),
             ("     골격 통제 ", 16, GREY),
             (f"{T['partial_spearman_controlling_tanimoto']:+.3f}", 19, ORANGE, True),
             ("     골격+크기 ", 16, GREY),
             (f"{T['partial_spearman_controlling_both']:+.3f}", 19, ORANGE, True)]],
           align=PP_ALIGN.CENTER)

    # ── 10 near 대역 ────────────────────────────────────────────────
    s = d.slide("한 대역만 유의하지만, 그 신호는 칸 하나에 얹혀 있다", "결과 ②",
                "near 대역, 즉 실데나필과 닮은 화합물들만 유의합니다. 보정 후에도 살아남고 "
                "교란 통제하면 오히려 강해집니다. 그런데 세 가지 때문에 무게를 실을 수 없습니다. "
                "AUC 신뢰구간이 기준 0.7을 포함하고, 그 대역의 약함 칸이 전수라서 빼면 "
                "유의하지 않고, 그 칸 역가가 어세이 바닥값에 몰려 있습니다. 게다가 그 대역조차 "
                "자세가 25%만 기준 안에 있습니다.\n"
                "[다음으로] 여기까지가 이번 데이터입니다. 남은 질문이 하나 있습니다.")
    for i, (v, lab, c) in enumerate([
            (f"{wb['near']['spearman']:+.3f}", f"near 대역 ρ   (p = {wb['near']['perm_p']})", GREEN),
            (f"{ACI['near']['auc']}",
             f"AUC   95% CI [{ACI['near']['ci95'][0]:.2f}, {ACI['near']['ci95'][1]:.2f}]", ORANGE),
            (f"{NSEN['drop_weak_cell']['spearman']:+.3f}", "전수 칸 제외 시 ρ", ORANGE)]):
        d.bignum(s, 16 + i * ((W - 32) / 3), 40, (W - 32) / 3, v, lab, c, sz=44)
    d.box(s, 22, 78, W - 44, 34, fill=RGBColor(0xFD, 0xF0, 0xE6), line=ORANGE, lw=1.6)
    d.text(s, 28, 82, W - 56, 28,
           [[("이 결과에 무게를 실을 수 없다", 19, ORANGE, True)],
            [("(a) AUC 신뢰구간이 사전 기준 0.7 을 포함한다   "
              "(b) 전수(census)인 약함 칸을 빼면 유의하지 않다   "
              "(c) 그 칸의 역가가 어세이 바닥값에 몰려 있다", 15, INK)]], space=4)
    d.box(s, 22, 118, W - 44, 40, fill=LGREY, line=NAVY, lw=1.4)
    d.text(s, 28, 122, W - 56, 34,
           [[("그리고 대역별 자세 타당도가 해석을 다시 제한한다", 17, NAVY, True)],
            [(f"mid {PV['mid']['frac_under_threshold']:.1%}   ·   "
              f"near {PV['near']['frac_under_threshold']:.1%} 만 2 Å 기준 안에 있다", 16, INK)],
            [("far 대역은 공유 부분구조가 평균 10.6원자(하한 8)라 이 지표 자체가 성립하지 "
              "않아 값을 내지 않았다. 이 지표는 “자세가 맞는가” 가 아니라 “공유 부분구조가 "
              "공결정 리간드 위에 겹치는가” 를 잰다.", 14, GREY)]],
           space=3)

    # ── 11 별개 질문: 비재현 ────────────────────────────────────────
    s = d.slide("다른 질문 — n=30 의 −0.538 은 어디서 왔나 (우리 초안의 답도 틀렸다)",
                "별개 질문 — 비재현",
                "여기서 주제가 바뀝니다. 앞은 자세와 상관 이야기였고, 이건 우리 자신에 대한 "
                "이야기입니다. 같은 하네스가 n=30에서 -0.538을 냈었습니다. 이 판본 초안은 그게 "
                "골격 교란 때문이라고 적었습니다. 그런데 옛 데이터에 같은 통제를 적용해 보니 "
                "1.2%만 줄었습니다. 설계 재현도, 프로토콜 재도킹도 원인이 아니었습니다. "
                "원인은 그 화합물 집합 자체이고, 어떤 성질 때문인지는 특정하지 못했습니다.\n"
                "[다음으로] 정리하겠습니다.")
    tests = [("① 옛 데이터에 골격 통제",
              f"{col['old_design']['spearman_top_pose']:+.3f} → "
              f"{col['old_design']['partial_controlling_tanimoto']:+.3f}",
              f"감쇠 {col['old_design']['attenuation_from_scaffold_control']:.1%}"),
             ("② 옛 데이터의 점수 vs Tanimoto",
              f"{col['old_design']['spearman_score_vs_tanimoto']:+.3f}",
              f"p = {col['old_design']['perm_p_score_vs_tanimoto']}"),
             ("③ 옛 설계를 새 데이터에 재현",
              f"중앙값 {col['test_1_design']['median']:+.3f}",
              f"−0.4 미만 {col['test_1_design']['frac_below_-0.4']:.1%}"),
             ("④ 옛 30건을 새 프로토콜로 재도킹",
              f"{oldp['rho_old_protocol']:+.3f} → {oldp['rho_new_protocol']:+.3f}" if oldp else "—",
              f"점수 상관 {oldp['spearman_between_protocols']:+.3f}" if oldp else "—")]
    for i, (nm, v1, v2) in enumerate(tests):
        yy = 38 + i * 19
        d.box(s, 20, yy, W - 40, 17, fill=(LGREY if i % 2 == 0 else WHITE),
              shape=MSO_SHAPE.RECTANGLE)
        d.text(s, 25, yy + 2, 130, 13, [(nm, 16, INK)], anchor=MSO_ANCHOR.MIDDLE)
        d.text(s, 158, yy + 2, 82, 13, [(v1, 18, ORANGE, True)],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        d.text(s, 244, yy + 2, 74, 13, [(v2, 15, GREY)],
               align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    d.box(s, 20, 118, W - 40, 40, fill=LGREY, line=NAVY, lw=1.6)
    d.text(s, 26, 122, W - 52, 34,
           [[("이 판본의 초안은 “골격 교란 때문” 이라고 적었다. 검정해 보니 1.2% 였다.",
              18, NAVY, True)],
            [("네 검정이 모두 후보를 배제한다. 원인은 그 화합물 집합 자체이며, 어떤 성질 "
              "때문인지는 이 데이터로 더 좁히지 못했다 — 그렇게 적었다.", 15, INK)]], space=4)

    # ── 12 정리 ─────────────────────────────────────────────────────
    s = d.slide("무엇을 말할 수 있고, 무엇이 남았나", "정리",
                "앞에서 말씀드린 결과는 반복하지 않겠습니다. 남은 것과 가져가실 것만 "
                "말씀드립니다. 배제한 것과 배제하지 못한 것을 구분해서 적었습니다.\n"
                "[다음으로] 마지막으로 이 연구가 자기 자신에 대해 알아낸 것을 말씀드리겠습니다.")
    d.box(s, 20, 38, (W - 46) / 2, 56, fill=LGREY, line=GREEN, lw=1.4)
    d.text(s, 26, 42, (W - 46) / 2 - 12, 50,
           [[("실험으로 제거한 것", 18, GREEN, True)],
            [("C2 자세순위 실패의 원인 2개", 15, INK, True)],
            [("   탐색 깊이 · 프로토네이션", 14, GREY)],
            [("n=30 비재현의 설명 2개", 15, INK, True)],
            [("   골격 교란 · 프로토콜 차이", 14, GREY)]], space=2.5)
    d.box(s, 26 + (W - 46) / 2, 38, (W - 46) / 2, 56, fill=LGREY, line=ORANGE, lw=1.4)
    d.text(s, 32 + (W - 46) / 2, 42, (W - 46) / 2 - 12, 50,
           [[("아직 배제하지 못한 것", 18, ORANGE, True)],
            [("물 전량 제거 · 강체 수용체", 15, INK)],
            [("단일 배좌 · holo 유도적합 편향", 15, INK)],
            [("그리고 채점 함수 자체", 15, INK, True)],
            [("   이들은 검정하지 않았다", 14, GREY)]], space=2.5)
    d.text(s, 20, 98, W - 40, 10,
           [("두 갈래는 서로 다른 현상이다 — 하나로 세지 않는다", 15, NAVY, True)],
           align=PP_ALIGN.CENTER)
    d.box(s, 20, 112, W - 40, 46, fill=RGBColor(0xE8, 0xF4, 0xEF), line=GREEN, lw=1.8)
    d.text(s, 26, 116, W - 52, 40,
           [[("가져가실 것 셋", 19, GREEN, True)],
            [("①  재도킹 대조를 C1/C2 로 쪼개고 깊이를 흔들어 병목을 특정하라", 16, INK)],
            [(f"②  점수를 매긴 자세가 어디에 놓였는지 먼저 보고하라 — 여기서는 "
              f"{PVI['frac_under_threshold']:.0%} 였다", 16, INK)],
            [("③  벤치마크의 골격 분포를 통제하라 (다만 여기서 그 통제가 바꾼 몫은 거의 0 이었다)",
              16, INK)]], space=3)

    # ── 13 마무리 ───────────────────────────────────────────────────
    s = d.prs.slides.add_slide(d.blank); d.n += 1
    d.box(s, 0, 0, W, H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    d.text(s, 24, 46, W - 48, 24,
           [("같은 데이터로 세 판본, 세 결론.\n매번 게이트는 전부 통과했다.", 32, WHITE, True)],
           align=PP_ALIGN.CENTER)
    d.text(s, 24, 92, W - 48, 20,
           [("방금 보신 11번 슬라이드가 그 세 번째입니다 — 이 판본의 초안도 틀렸고, "
             "그걸 잡은 것은 자동 검증이 아니라 외부 비평이었습니다.", 16, PALE)],
           align=PP_ALIGN.CENTER)
    d.text(s, 24, 118, W - 48, 14,
           [("게이트는 수치가 산출되었는지를 검사할 뿐, 결론이 데이터에서 따라 나오는지는 "
             "검사하지 못한다.", 17, WHITE, True)], align=PP_ALIGN.CENTER)
    d.text(s, 24, 146, W - 48, 16,
           [[("세 판본의 차이와 리뷰 기록 전부 공개", 14, PALE)],
            [("github.com/fourmodern/2026_aidrugdiscovery → Day06_LLM_Agent/pde5_harness",
              16, WHITE, True)]], align=PP_ALIGN.CENTER, space=3)
    s.notes_slide.notes_text_frame.text = (
        "질문 대비: (1) 왜 PDE5A 인가 — 결정 구조와 활성 데이터가 모두 풍부해서. "
        "(2) 다른 표적에서도 같은가 — 검정하지 않았고 그것이 가장 큰 한계. 표적 하나, 구조 "
        "하나입니다. (3) 물을 넣으면 달라지나 — 배제하지 못한 후보이며 다음 실험. "
        "(4) 자세가 81% 틀렸는데 상관을 논할 수 있나 — 그래서 6번에서 미리 상한을 걸었고, "
        "결론도 그 안에서만 말했습니다.")

    out = Path(a.out); out.parent.mkdir(parents=True, exist_ok=True)
    n = d.save(out)
    print(f"{out}  ({n}장, 16:9)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
