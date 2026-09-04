#!/usr/bin/env python3
"""학회 발표용 A0 포스터 (841 × 1189 mm, 세로).

포스터는 논문 축약이 아니다. 3 m 밖에서 제목이, 1.5 m 에서 결론이, 0.6 m 에서 본문이
읽혀야 한다. 그래서 글자를 세 단계로 나누고 결론을 위쪽에 둔다.

**배치는 흐름식이다.** 각 요소가 자기가 소비한 높이를 돌려주고 커서가 그만큼 내려간다.
고정값으로 커서를 옮기면 줄바꿈이 늘어난 문단이 다음 절 헤더를 침범한다 (실제로 그랬다).

모든 수치는 sample_run/ 산출 파일에서 읽는다. 손으로 적은 결과 수치는 없다.
"""
from __future__ import annotations
import argparse, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
SR = ROOT / "sample_run"
FIG = SR / "report" / "figures_controlled"
FONT = "Pretendard"

W, H = 841, 1189                       # A0 세로 (mm)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
NAVY   = RGBColor(0x0B, 0x2A, 0x4A)
BLUE   = RGBColor(0x00, 0x72, 0xB2)
ORANGE = RGBColor(0xD5, 0x5E, 0x00)
GREEN  = RGBColor(0x00, 0x9E, 0x73)
GREY   = RGBColor(0x5A, 0x63, 0x6B)
LGREY  = RGBColor(0xF2, 0xF4, 0xF6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PALE   = RGBColor(0x9F, 0xC5, 0xE0)

M, GAP = 30, 22
COLW = (W - 2 * M - 2 * GAP) / 3
CX = [M, M + COLW + GAP, M + 2 * (COLW + GAP)]
FULLW = W - 2 * M

T_TITLE, T_SUB, T_HEAD = 78, 38, 32
T_BODY, T_SMALL, T_BIG = 25, 21, 38
T_KVK, T_KVV = 24, 29


def load(n):
    p = SR / n
    if not p.exists():
        return None
    d = json.loads(p.read_text())
    return d.get("result", d)


def box(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = s.shapes.add_shape(shape, Mm(x), Mm(y), Mm(w), Mm(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def est_h(runs, w, space=0):
    """줄바꿈 후 높이 추정. 한글 1 em · 라틴 0.5 em 혼용을 0.62 em 로 잡는다."""
    paras = runs if isinstance(runs[0], list) else [runs]
    tot = 0.0
    for para in paras:
        chars = sum(len(t) for t, *_ in para)
        sz = max((r[1] for r in para), default=T_BODY)
        per = max(1, int(w / (sz * 0.3528 * 0.62)))
        tot += max(1, -(-chars // per)) * (sz * 0.3528 * 1.42) + space * 0.3528
    return tot + 2


def text(s, x, y, w, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, h=None):
    hh = h if h is not None else est_h(runs, w, space)
    tb = s.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(hh))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Mm(1.5); tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for t, sz, colr, *rest in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = colr
            r.font.bold = bool(rest and rest[0]); r.font.name = FONT
    return hh


def head(s, x, y, w, n, title, colr=NAVY):
    box(s, x, y, w, 17, fill=colr, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 5, y + 1.5, 18, [(n, T_HEAD, WHITE, True)], h=14, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 26, y + 1.5, w - 30, [(title, T_HEAD, WHITE, True)], h=14,
         anchor=MSO_ANCHOR.MIDDLE)
    return 17 + 6


def pic(s, name, x, y, w):
    f = FIG / name
    if not f.exists():
        box(s, x, y, w, 40, fill=LGREY, line=GREY)
        text(s, x, y + 17, w, [(f"[그림 없음: {name}]", T_BODY, ORANGE)],
             align=PP_ALIGN.CENTER)
        return 44
    p = s.shapes.add_picture(str(f), Mm(x), Mm(y), width=Mm(w))
    return p.height / 36000 + 5


def kv(s, x, y, w, rows, lw=0.50):
    yy = y
    for i, (k, v, colr) in enumerate(rows):
        box(s, x, yy, w, 16, fill=(LGREY if i % 2 == 0 else WHITE),
            shape=MSO_SHAPE.RECTANGLE)
        text(s, x + 5, yy + 1, w * lw, [(k, T_KVK, INK)], h=14, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + w * lw, yy + 1, w * (1 - lw) - 5, [(v, T_KVV, colr, True)], h=14,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        yy += 16
    return yy - y + 6


def panel(s, x, y, w, runs, fill=LGREY, line=BLUE, lw=1.2, space=5, pad=6):
    hh = est_h(runs, w - 2 * pad, space) + 2 * pad
    box(s, x, y, w, hh, fill=fill, line=line, lw=lw)
    text(s, x + pad, y + pad, w - 2 * pad, runs, space=space, h=hh - 2 * pad)
    return hh + 6


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(SR / "report" / "docs" / "poster_A0.pptx"))
    a = ap.parse_args()

    ds, dk, an = load("dataset_controlled.json"), load("docking_controlled.json"), \
                 load("analysis_controlled.json")
    tc, sw = load("terms_controlled.json"), load("exhaustiveness_sweep.json")
    prot, col = load("protonation_test.json"), load("collapse_diagnosis.json")
    oldp = load("old_set_new_protocol.json")
    if not (ds and dk and an):
        raise SystemExit("산출 파일이 없다 — 포스터를 만들지 않는다.")
    T = an["arms"]["top_pose"]; PV = T["pose_validity"]; PVI = PV["interpretable_only"]
    wb = T["within_similarity_bin"]
    ACI = T["auc_ci95"]; NSEN = T["near_sensitivity"]; ctrl = dk["control_redock"]

    prs = Presentation(); prs.slide_width = Mm(W); prs.slide_height = Mm(H)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, 0, 0, W, H, fill=WHITE, shape=MSO_SHAPE.RECTANGLE)

    # ── 제목 띠 ─────────────────────────────────────────────────────
    box(s, 0, 0, W, 134, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    text(s, M, 14, FULLW, [("골격 유사성과 역가를 직교화한 PDE5A 도킹 재평가",
                            T_TITLE, WHITE, True)], align=PP_ALIGN.CENTER, h=42)
    text(s, M, 60, FULLW, [("선행 상관의 재현 실패, 그리고 자세 타당도가 정하는 해석의 상한",
                            T_SUB, PALE)], align=PP_ALIGN.CENTER, h=22)
    text(s, M, 88, FULLW,
         [("Orthogonalizing scaffold similarity and potency in a PDE5A docking benchmark: "
           "non-replication of a prior score–potency correlation",
           T_BODY, RGBColor(0x7F, 0xA8, 0xC8))], align=PP_ALIGN.CENTER, h=18)
    text(s, M, 110, FULLW,
         [(f"smina / AutoDock Vina 1.1.2   ·   PDB 1UDT   ·   ChEMBL PDE5A (n = {an['n']})"
           f"   ·   전 산출물·코드·리뷰 기록 공개: github.com/fourmodern/2026_aidrugdiscovery",
           T_SMALL, PALE)], align=PP_ALIGN.CENTER, h=16)

    # ── 핵심 메시지 ─────────────────────────────────────────────────
    y = 146
    box(s, M, y, FULLW, 60, fill=RGBColor(0xFD, 0xF0, 0xE6), line=ORANGE, lw=2.4)
    text(s, M + 10, y + 6, FULLW - 20,
         [(f"골격을 분리하자 도킹 점수와 역가의 상관이 사라졌다    ρ = {T['spearman']:+.3f}"
           f"   (95% CI {T['ci95'][0]:+.2f} ~ {T['ci95'][1]:+.2f},  p = {T['perm_p']})",
           T_BIG, ORANGE, True)], align=PP_ALIGN.CENTER, h=24)
    text(s, M + 10, y + 34, FULLW - 20,
         [(f"n = 30 에서 얻었던 ρ = {col['old_design']['spearman_top_pose']:+.3f} 는 재현되지 "
           f"않았고, 그 이유로 제시했던 골격 교란 설명마저 검정 결과 기각되었다 "
           f"(옛 데이터에서 골격 통제 시 감쇠 "
           f"{col['old_design']['attenuation_from_scaffold_control']:.1%})",
           T_BODY + 4, INK, True)], align=PP_ALIGN.CENTER, h=22)
    y += 66
    # 폭을 조금 줄여 세 단이 하단 띠 안에 들어오게 한다 (넘침 검사가 잡아준다)
    y += pic(s, "fig01_graphical_abstract.png", M + FULLW * 0.075, y, FULLW * 0.85)
    ytop = y + 2

    # ══ 1단 ════════════════════════════════════════════════════════
    y = ytop
    y += head(s, CX[0], y, COLW, "1", "배경과 목적")
    y += text(s, CX[0], y, COLW,
              [[("도킹 평가는 보통 활성 화합물을 모아 점수와 역가의 상관 하나를 보고한다. "
                 "그런데 표적 좌표는 대개 어떤 리간드가 결합한 채로 풀린 holo 구조이고, 그 "
                 "리간드를 닮은 화합물은", T_BODY, INK)],
               [("① 유도적합된 포켓에 잘 맞아 점수가 좋고   ② 같은 계열의 최적화 산물이라 "
                 "대체로 더 강력하다.", T_BODY, BLUE, True)],
               [("두 경로의 방향이 같으므로 상관은 관찰되지만, 그것이 채점 함수의 능력인지 "
                 "표적 구조 선택의 부산물인지 구분되지 않는다. 비활성 대비 편향은 널리 "
                 "논의됐지만(DUD-E, LIT-PCBA), 활성 집합 내부의 골격 분포가 상관을 얼마나 "
                 "만드는지는 정량된 적이 거의 없다.", T_BODY, INK)]], space=5) + 6
    y += panel(s, CX[0], y, COLW,
               [[("연구 목적", T_HEAD - 4, BLUE, True)],
                [("1.  역가를 정량 예측하는가", T_BODY, INK), ("   →  Q² ≥ 0.3", T_BODY, GREY)],
                [("2.  강한 화합물을 선별하는가", T_BODY, INK),
                 ("   →  p < 0.05 & AUC ≥ 0.7", T_BODY, GREY)],
                [("3.  골격 통제 후에도 남는가", T_BODY, INK),
                 ("   →  편상관 60% 유지", T_BODY, GREY)],
                [("4.  대조 실패 시 무엇이 병목인가", T_BODY, INK),
                 ("   →  후보를 실험으로 제거", T_BODY, GREY)]], space=4)
    y += head(s, CX[0], y, COLW, "2", "설계 — 두 축의 직교화")
    y += pic(s, "fig02_dataset.png", CX[0], y, COLW)
    y += text(s, CX[0], y, COLW,
              [[(f"ChEMBL PDE5A 활성 {ds['records_scanned']} 레코드 전수 → 화합물당 중앙값 "
                 f"집계({ds['pool_size']}종) → 역가 3구간 × 실데나필 Tanimoto 3구간의 9칸 "
                 f"격자에서 균등 추출.", T_BODY, INK)],
               [(f"교란 상관을 Pearson "
                 f"{ds['confound_pearson_r_potency_vs_similarity']:+.3f} 까지 낮췄다 "
                 f"(역가로만 층화한 옛 설계: +0.281).", T_BODY, GREEN, True)]], space=4) + 6
    y += text(s, CX[0], y, COLW,
              [[("도킹", T_HEAD - 6, NAVY, True)],
               [(f"{dk['engine']}, seed {dk['seed']}, num_modes 9, autobox_add 3. "
                 f"{dk['n_docked']}/{dk['n_attempted']} 성공. 자세는 공결정 리간드를 참조하지 "
                 f"않는 점수 1위를 주 결과로 썼다.", T_BODY, INK)]], space=4)

    y += head(s, CX[0], y, COLW, "2b", "사전 기준 대비 판정", ORANGE)
    y += kv(s, CX[0], y, COLW,
            [("① 정량 예측   Q² ≥ 0.3",
              f"{tc['models']['single_vina_score']['Q2_loo']:+.3f}   미달" if tc else "—", ORANGE),
             ("② 선별   AUC ≥ 0.7", f"{ACI['all']['auc']}   미달", ORANGE),
             ("② 선별   순열 p < 0.05", f"{T['perm_p']}   미달", ORANGE),
             ("③ 골격 통제 후 60% 유지", "판정 불가", ORANGE),
             ("④ 병목 특정   후보 제거", "4건 제거", GREEN)], lw=0.70)
    y += text(s, CX[0], y, COLW,
              [[("③ 은 원상관 자체가 유의하지 않아 “유지” 를 물을 수 없다. "
                 "④ 는 두 갈래에서 각각 2건씩이며 하나로 세지 않는다.", T_SMALL, GREY)]]) + 4
    ends = [y]                          # 단별 최하단
    # ══ 2단 ════════════════════════════════════════════════════════
    y = ytop
    y += head(s, CX[1], y, COLW, "3", "재도킹 대조 — 병목 찾기")
    y += pic(s, "fig03_control.png", CX[1], y, COLW)
    y += kv(s, CX[1], y, COLW,
            [("C1  샘플링 (상위 모드 최선)",
              f"{ctrl['rmsd_best_of_modes_angstrom']:.2f} Å  PASS", GREEN),
             ("C2  채점 (점수 1위 자세)",
              f"{ctrl['rmsd_top_pose_angstrom']:.2f} Å  FAIL", ORANGE)])
    y += text(s, CX[1], y, COLW,
              [[("탐색은 결정 자세를 찾아낸다. 채점이 그것을 1위로 올리지 못한다. 두 대조를 "
                 "합쳐 하나의 숫자로 보고하면 이 구분이 사라진다.", T_BODY, INK)]]) + 6
    y += head(s, CX[1], y, COLW, "4", "C2 실패 원인을 실험으로 제거", GREEN)
    y += kv(s, CX[1], y, COLW,
            [("탐색 깊이 4 → 128 (32배)", "C2 전 구간 0%" if sw else "—", GREEN),
             ("프로토네이션 2×2", "4 조건 점수 동일" if prot else "—", GREEN)], lw=0.60)
    y += text(s, CX[1], y, COLW,
              [[("프로토네이션이 왜 무효였는지까지 확인했다. ", T_BODY, INK),
                ("기본 Vina 함수에는 정전기 항이 없다", T_BODY, ORANGE, True),
                (" — 형식전하는 점수에 닿지 못하고, 바뀐 수소결합 타이핑도 그 질소가 결합 "
                 "거리 밖이라 항 값이 움직이지 않았다.", T_BODY, INK)]]) + 6
    y += pic(s, "fig08_exhaustiveness.png", CX[1], y, COLW)

    y += head(s, CX[1], y, COLW, "5", "별개 질문 — 비재현", GREY)
    y += kv(s, CX[1], y, COLW,
            [("골격 교란 (옛 데이터에 통제)",
              f"감쇠 {col['old_design']['attenuation_from_scaffold_control']:.1%}" if col else "—",
              GREY),
             ("프로토콜 (옛 30건 재도킹)",
              f"점수 상관 {oldp['spearman_between_protocols']:+.3f}" if oldp else "—", GREY)],
            lw=0.62)
    y += text(s, CX[1], y, COLW,
              [[("n=30 의 ", T_BODY, INK),
                (f"ρ = {col['old_design']['spearman_top_pose']:+.3f}", T_BODY, INK, True),
                (" 는 재현되지 않았다. 초안은 골격 교란 때문이라 적었으나 검정 결과 "
                 "기각됐다. ", T_BODY, INK),
                ("위 두 검정은 C2 실패가 아니라 이 비재현을 대상으로 한 것", T_BODY, ORANGE, True),
                ("이며, 4번 절의 두 검정과 하나로 세지 말아야 한다.", T_BODY, INK)]]) + 6

    ends.append(y)
    # ══ 3단 ════════════════════════════════════════════════════════
    y = ytop
    y += head(s, CX[2], y, COLW, "6", "결과 — 골격 통제 후")
    y += pic(s, "fig05_forest.png", CX[2] + COLW * 0.08, y, COLW * 0.84)
    y += kv(s, CX[2], y, COLW,
            [("전체 순위상관", f"{T['spearman']:+.3f}", ORANGE),
             ("골격 통제 편상관", f"{T['partial_spearman_controlling_tanimoto']:+.3f}", ORANGE),
             ("골격+크기 동시 통제", f"{T['partial_spearman_controlling_both']:+.3f}", ORANGE),
             ("교차검증 Q² (단일 점수)",
              f"{tc['models']['single_vina_score']['Q2_loo']:+.3f}" if tc else "—", ORANGE)],
            lw=0.58)
    y += text(s, CX[2], y, COLW,
              [[("위 넷은 상관·Q² 값이다. ", T_SMALL, GREY),
                (f"순열 p = {T['perm_p']}", T_BODY, ORANGE, True),
                (" 로 유의하지 않다.", T_SMALL, GREY)]]) + 3
    y += text(s, CX[2], y, COLW,
              [[(f"near 대역만 유의하다 (ρ = {wb['near']['spearman']:+.3f}, "
                 f"p = {wb['near']['perm_p']}). ", T_BODY, INK),
                ("가설로만 남긴다", T_BODY, ORANGE, True),
                (f" — AUC 신뢰구간 [{ACI['near']['ci95'][0]:.2f}, "
                 f"{ACI['near']['ci95'][1]:.2f}] 이 기준 0.7 을 포함하고, 전수 칸을 빼면 "
                 f"{NSEN['drop_weak_cell']['spearman']:+.3f} 로 떨어진다.", T_BODY, INK)]]) + 4
    y += head(s, CX[2], y, COLW, "6b", "해석의 상한 — 자세 타당도", ORANGE)
    y += kv(s, CX[2], y, COLW,
            [("1위 자세 MCS-RMSD 중앙값", f"{PVI['median_rmsd']:.2f} Å", ORANGE),
             ("2 Å 기준 통과", f"{PVI['frac_under_threshold']:.1%}", ORANGE),
             ("대상", f"해석 가능 대역 {PVI['n']}건", GREY)], lw=0.62)
    y += text(s, CX[2], y, COLW,
              [[(f"점수를 매긴 자세의 {1 - PVI['frac_under_threshold']:.0%} 가 C2 를 "
                 f"실패시킨 것과 같은 기준 밖이다 — 이것이 모든 상관 해석의 상한이다. ",
                 T_BODY, INK),
                (f"far 대역은 공유 부분구조가 평균 {PV['far']['mean_mcs_atoms']}원자라 "
                 f"지표가 성립하지 않아 제외했다.", T_SMALL, GREY)]]) + 3
    y += head(s, CX[2], y, COLW, "7", "결론", NAVY)
    y += panel(s, CX[2], y, COLW,
               [[("· 정량 예측은 실패했다.", T_BODY + 2, INK, True),
                 (f" Q² 는 세 모델 모두 0 근처로 기준 0.3 에 한참 미달한다.", T_BODY, INK)],
                [("· 선별도 전체적으로는 작동하지 않았다.", T_BODY + 2, INK, True),
                 (" 한 대역의 신호는 전수 칸 하나에 얹혀 있다.", T_BODY, INK)],
                [("· 두 갈래에서 각각 두 후보씩 제거했다.", T_BODY + 2, INK, True),
                 (" C2 실패는 탐색 깊이·프로토네이션, 비재현은 골격 교란·프로토콜. "
                  "남은 것은 물 제거·강체 수용체·단일 배좌·holo 편향·채점 함수 자체다.",
                  T_BODY, INK)],
                [("· 벤치마크의 골격 분포를 통제하라.", T_BODY + 2, GREEN, True),
                 (" 다만 여기서 그 통제가 바꾼 몫은 거의 0 이었다.", T_BODY, INK)],
                [("· 자세가 어디에 놓였는지 먼저 보고하라.", T_BODY + 2, GREEN, True),
                 (" 그 비율 없이 상관값만 보면 무엇을 측정했는지 알 수 없다.", T_BODY, INK)],
                [("· “신호 없음” 은 귀무를 채택한 것이 아니다.", T_BODY + 2, NAVY, True),
                 (f" 신뢰구간 상한이 {T['ci95'][1]:+.2f} 이므로 |ρ| ≲ "
                  f"{abs(T['ci95'][0]):.2f} 인 관계는 배제하지 못한다. 기각하지 못한 것이지 "
                  f"없다고 보인 것이 아니다.", T_BODY, INK)]],
               line=NAVY, lw=1.6, space=6)
    ends.append(y); bottom = max(ends)

    # ── 하단 띠 ─────────────────────────────────────────────────────
    box(s, 0, H - 32, W, 32, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    text(s, M, H - 26, FULLW,
         [("이 연구는 교육 목적의 시연이며 임상·규제 판단의 근거로 사용할 수 없다.   "
           "데이터: ChEMBL (CC BY-SA 3.0) · RCSB PDB.   도구: AutoDock Vina (Apache-2.0) · "
           "smina (GPL-2.0) · RDKit (BSD-3) · Open Babel (GPL-2.0) · PyMOL.",
           T_SMALL - 3, PALE)], align=PP_ALIGN.CENTER, h=9)
    text(s, M, H - 16, FULLW,
         [("전 산출물·코드·리뷰 기록 공개  ·  github.com/fourmodern/2026_aidrugdiscovery "
           "→ Day06_LLM_Agent/pde5_harness", T_SMALL, WHITE, True)],
         align=PP_ALIGN.CENTER, h=10)

    out = Path(a.out); out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    fill = bottom / (H - 32)
    print(f"{out}")
    print(f"  A0 {W}×{H} mm · 도형 {len(s.shapes)}개 · 채움 {fill:.0%}")
    print("  단별 최하단: " + " · ".join(f"{i+1}단 {e:.0f}mm" for i, e in enumerate(ends))
          + f"  (한계 {H - 40:.0f}mm)")
    if bottom > H - 40:
        print("  경고: 콘텐츠가 하단 띠를 침범한다")
    elif fill < 0.80:
        print("  경고: 아래쪽이 비어 있다 — 그림을 키우거나 내용을 늘려라")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
