#!/usr/bin/env python3
"""배포 전 자체 감사 — 보고서의 주장과 산출 파일이 실제로 맞는지 기계로 확인한다.

리뷰에서 반복해서 걸린 실패 유형이 있다. 본문을 고쳤는데 그림·슬라이드·PDF 가 옛 주장을
그대로 담고 있는 것. 사람이 눈으로 확인하겠다는 약속은 세 번 다 실패했으므로 코드로 만든다.
"""
from __future__ import annotations
import hashlib, json, re, sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
SR = ROOT / "sample_run"
RP = SR / "report"


def load(n):
    p = SR / n
    if not p.exists(): return None
    d = json.loads(p.read_text())
    return d.get("result", d)


def env(n):
    p = SR / n
    return json.loads(p.read_text()) if p.exists() else None


def main() -> int:
    fails, warns = [], []
    md = RP / "report_controlled.md"
    if not md.exists():
        print("report_controlled.md 가 없다 — 감사 불가"); return 1
    text = md.read_text()

    # 1. 모든 게이트가 PASS 인가
    arts = ["dataset_controlled.json", "docking_controlled.json", "exhaustiveness_sweep.json",
            "analysis_controlled.json", "terms_controlled.json", "statistics_validation.json",
            "contact_concordance.json", "custom_scoring_controlled.json"]
    for a in arts:
        e = env(a)
        if e is None:
            fails.append(f"산출 파일 없음: {a}"); continue
        v = e.get("verification")
        if not v:
            fails.append(f"검증 봉투 없음: {a}")
        elif not v.get("passed"):
            bad = [c["check"] for c in v.get("checks", []) if not c.get("passed")]
            fails.append(f"게이트 FAIL: {a} — {bad}")

    # 2. 본문이 참조하는 그림이 실제로 존재하는가
    imgs = re.findall(r"!\[[^\]]*\]\(([^)]+)\)", text)
    missing = [i for i in imgs if not (md.parent / i).exists()]
    if missing:
        fails.append(f"본문이 참조하는 그림 누락: {missing}")

    # 3. 그림이 본문보다 오래되지 않았는가 (본문 수정 후 그림 미재생성 탐지)
    mt = md.stat().st_mtime
    stale = [i for i in imgs if (md.parent / i).exists()
             and (md.parent / i).stat().st_mtime < mt - 300]
    if stale:
        warns.append(f"본문보다 5분 이상 오래된 그림: {stale}")

    # 4. 배포 문서가 본문·그림보다 새로운가
    # 배포물 전체를 본다. 이름 하나만 검사하면 포스터·발표덱이 사각지대가 된다
    # (실제로 pptx 표지가 옛 제목을 달고 있는 것을 사용자가 먼저 발견했다).
    docs = sorted(f for f in (RP / "docs").glob("*")
                  if f.suffix in (".docx", ".pptx", ".pdf")) if (RP / "docs").exists() else []
    # 배포물은 셋뿐이다: 문서(docx/pdf) · 학회 포스터(A0) · 구두발표 덱(16:9).
    # 보고서에서 자동 파생하던 슬라이드는 발표 덱과 목적이 겹치고 핵심 수치도
    # 싣지 못해 제거했다.
    EXPECT = {"report_controlled.docx", "report_controlled_doc.pdf",
              "poster_A0.pptx", "poster_A0.pdf", "talk_16x9.pptx", "talk_16x9.pdf"}
    missing = EXPECT - {f.name for f in docs}
    if missing:
        fails.append(f"배포물 누락: {sorted(missing)}")
    base = max([mt] + [(md.parent / i).stat().st_mtime for i in imgs if (md.parent / i).exists()])
    old = [d.name for d in docs if d.stat().st_mtime < base]
    if not docs:
        fails.append("배포 문서가 없다")
    elif old:
        fails.append(f"본문·그림보다 오래된 배포 문서: {old}")

    # 5. 본문의 SHA256 표기가 실제 파일과 맞는가
    for name, tag in (("dataset_controlled.json", "dataset"),
                      ("docking_controlled.json", "docking"),
                      ("analysis_controlled.json", "analysis")):
        m = re.search(rf"{tag} `([0-9a-f]{{16}})`", text)
        if not m:
            warns.append(f"본문에 {tag} 해시 표기 없음"); continue
        real = hashlib.sha256((SR / name).read_bytes()).hexdigest()[:16]
        if m.group(1) != real:
            fails.append(f"해시 불일치 {name}: 본문 {m.group(1)} vs 실제 {real}")

    an = load("analysis_controlled.json")

    # 6. 철회한 주장이 남아 있지 않은가
    retracted = ["신뢰할 수준으로 재현하지 못했다",
                 "결합 양상은 문헌과 일치한다",
                 "도킹이 화학적으로 말이 되는 위치"]
    for r in retracted:
        if r in text:
            fails.append(f"철회한 문장이 본문에 남아 있다: {r!r}")

    # 7. 핵심 수치가 산출 파일과 일치하는가
    if an:
        T = an["arms"]["top_pose"]
        for label, val in (("원상관", T["spearman"]),
                           ("편상관", T["partial_spearman_controlling_tanimoto"])):
            if f"{val:+.3f}" not in text:
                fails.append(f"{label} 값 {val:+.3f} 가 본문에 없다 — 재생성 필요")

    # 8. 무-날조: 금지 표현
    banned = ["predicted to be", "typical effect size", "plausible range",
              "illustrative", "expected shape", "대략 추정", "예상 효과크기"]
    hits = [b for b in banned if b in text]
    if hits:
        fails.append(f"금지 표현 검출: {hits}")

    # ── 배포 문서의 '내용'을 실제로 열어 본다 ────────────────────────
    # 여기까지 안 보면, 본문은 고쳤는데 사용자가 여는 파일은 옛 문구를 그대로 다는 일이
    # 계속 생긴다. 실제로 pptx 표지가 옛 하네스 제목을 달고 있었고 아무도 못 잡았다.
    def _docx_text(f):
        try:
            from docx import Document
            d = Document(str(f))
            return "\n".join(x.text for x in d.paragraphs)
        except Exception as e:
            return f"__ERR__{type(e).__name__}"

    def _pptx_text(f, first_slide_only=False):
        try:
            from pptx import Presentation
            pr = Presentation(str(f))
            sl = list(pr.slides)[:1] if first_slide_only else list(pr.slides)
            return "\n".join(sh.text_frame.text for s_ in sl for sh in s_.shapes
                              if sh.has_text_frame)
        except Exception as e:
            return f"__ERR__{type(e).__name__}"

    title = next((l.lstrip("# ").strip() for l in text.splitlines()
                  if l.startswith("# ")), None)
    for d in docs:
        body = (_docx_text(d) if d.suffix == ".docx"
                else _pptx_text(d) if d.suffix == ".pptx" else None)
        if body is None:
            continue
        if body.startswith("__ERR__"):
            warns.append(f"{d.name} 를 열지 못했다 ({body[7:]})"); continue
        # 제목 일치 — 표지가 다른 문서의 제목을 달고 있으면 독자가 오해한다
        if title and d.name.startswith("report_") and title[:18] not in body:
            fails.append(f"{d.name} 에 본문 제목('{title[:26]}…')이 없다 — 표지가 다른 "
                         f"문서의 것일 수 있다")
        # 철회된 문장이 배포본에만 남아 있는 경우
        for r in retracted:
            if r in body:
                fails.append(f"{d.name} 에 철회한 문장이 남아 있다: {r!r}")
        # 금지 표현
        for b in banned:
            if b in body:
                fails.append(f"{d.name} 에 금지 표현: {b!r}")
        if d.name.startswith(("poster", "talk")) and "PDE5A 도킹 재평가" not in body:
            fails.append(f"{d.name} 에 연구 제목이 없다")
        # 핵심 수치가 배포본에도 실려 있는가 (docx 한정 — 슬라이드는 발췌라 제외)
        if d.suffix in (".docx", ".pptx") and an and not d.name.endswith("_slides.pptx"):
            v = an["arms"]["top_pose"]["spearman"]
            if f"{v:+.3f}" not in body:
                fails.append(f"{d.name} 에 헤드라인 값 {v:+.3f} 가 없다")

    print("=" * 66)
    print(f"감사 대상  {md}")
    print(f"산출 파일  {len(arts)}개  ·  본문 참조 그림  {len(imgs)}장  ·  배포 문서 {len(docs)}개")
    print("=" * 66)
    for w in warns:
        print(f"  경고  {w}")
    for f in fails:
        print(f"  실패  {f}")
    if not fails:
        print("  전 항목 통과")
    print("=" * 66)
    return 1 if fails else 0


if __name__ == "__main__":
    raise SystemExit(main())
