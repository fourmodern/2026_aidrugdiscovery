#!/usr/bin/env python3
"""보고서 마크다운 → Word(.docx) + 발표용 PowerPoint(.pptx).

그림은 링크가 아니라 실제 이미지로, 표는 실제 표 객체로 삽입한다.
"""
import argparse, re
from pathlib import Path

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0B, 0x2A, 0x4A); CYAN = RGBColor(0x00, 0x9C, 0xD8)
GREY = RGBColor(0x5A, 0x63, 0x6B)
P_NAVY = PRGB(0x0B, 0x2A, 0x4A); P_CYAN = PRGB(0x00, 0x9C, 0xD8)
P_GREY = PRGB(0x5A, 0x63, 0x6B); P_WHITE = PRGB(0xFF, 0xFF, 0xFF)
KFONT = "Malgun Gothic"


def parse(md: str):
    """마크다운을 블록 목록으로 (heading / para / image / table / code)."""
    blocks, lines, i = [], md.split("\n"), 0
    while i < len(lines):
        ln = lines[i]
        if ln.startswith("!["):
            m = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", ln)
            if m: blocks.append(("image", m.group(2), m.group(1)))
            i += 1
        elif ln.startswith("```"):
            buf = []; i += 1
            while i < len(lines) and not lines[i].startswith("```"):
                buf.append(lines[i]); i += 1
            blocks.append(("code", "\n".join(buf), None)); i += 1
        elif ln.startswith("|") and i + 1 < len(lines) and re.match(r"^\|[\s:|-]+\|$", lines[i + 1]):
            rows = []
            while i < len(lines) and lines[i].startswith("|"):
                if not re.match(r"^\|[\s:|-]+\|$", lines[i]):
                    rows.append([c.strip() for c in lines[i].strip().strip("|").split("|")])
                i += 1
            blocks.append(("table", rows, None))
        elif ln.startswith("#"):
            lv = len(ln) - len(ln.lstrip("#"))
            blocks.append(("heading", ln.lstrip("#").strip(), lv)); i += 1
        elif ln.strip() in ("", "---"):
            i += 1
        else:
            buf = []
            while i < len(lines) and lines[i].strip() and not lines[i].startswith(("#", "|", "!", "```")):
                buf.append(lines[i]); i += 1
            blocks.append(("para", " ".join(buf), None))
    return blocks


def clean(t: str) -> str:
    return re.sub(r"\*\*(.+?)\*\*", r"\1", re.sub(r"`([^`]+)`", r"\1", t))


def build_docx(blocks, base: Path, out: Path):
    doc = Document()
    st = doc.styles["Normal"]; st.font.name = KFONT; st.font.size = Pt(10.5)
    n_img = n_tbl = 0
    for kind, val, extra in blocks:
        if kind == "heading":
            h = doc.add_heading(clean(val), level=min(extra, 4))
            for r in h.runs: r.font.color.rgb = NAVY; r.font.name = KFONT
        elif kind == "para":
            p = doc.add_paragraph()
            for seg in re.split(r"(\*\*.+?\*\*)", val):
                if not seg: continue
                r = p.add_run(clean(seg)); r.font.name = KFONT; r.font.size = Pt(10.5)
                r.bold = seg.startswith("**")
        elif kind == "image":
            src = (base / val).resolve()
            if src.exists():
                doc.add_picture(str(src), width=Inches(6.3)); n_img += 1
                doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif kind == "table":
            rows = val
            t = doc.add_table(rows=len(rows), cols=len(rows[0])); t.style = "Light Grid Accent 1"
            for ri, row in enumerate(rows):
                for ci, cell in enumerate(row[:len(rows[0])]):
                    c = t.cell(ri, ci); c.text = clean(cell)
                    for pr in c.paragraphs:
                        for r in pr.runs:
                            r.font.size = Pt(8.5); r.font.name = KFONT
                            r.bold = (ri == 0)
            n_tbl += 1
        elif kind == "code":
            p = doc.add_paragraph()
            r = p.add_run(val); r.font.name = "Consolas"; r.font.size = Pt(9)
            r.font.color.rgb = GREY
    doc.save(out)
    return n_img, n_tbl


def _tb(slide, x, y, w, h, runs, size=14, color=P_NAVY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(PIn(x), PIn(y), PIn(w), PIn(h)); tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = PPt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = KFONT
    return tb


def build_pptx(blocks, base: Path, out: Path):
    prs = Presentation(); prs.slide_width = PIn(13.333); prs.slide_height = PIn(7.5)
    BLANK = prs.slide_layouts[6]
    imgs = [(v, e) for k, v, e in blocks if k == "image"]
    heads = [(v, e) for k, v, e in blocks if k == "heading"]
    paras = [v for k, v, _ in blocks if k == "para"]

    def new(title, kicker=""):
        s = prs.slides.add_slide(BLANK)
        bar = s.shapes.add_shape(1, PIn(0), PIn(0), PIn(13.333), PIn(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = P_CYAN; bar.line.fill.background()
        if kicker: _tb(s, 0.6, 0.24, 11.8, 0.34, [kicker], size=11, color=P_CYAN, bold=True)
        _tb(s, 0.6, 0.58, 12.1, 0.9, [title], size=22, color=P_NAVY, bold=True)
        return s

    # 표지 — graphical abstract
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, PIn(0), PIn(0), PIn(13.333), PIn(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = P_NAVY; bg.line.fill.background()
    # 표지 문구는 본문 제목에서 가져온다 — 하드코딩하면 다른 보고서에 옛 제목이 붙는다
    hs = [v for k, v, lvl in blocks if k == "heading"]
    cover_title = hs[0] if hs else "연구 보고서"
    cover_sub = hs[1] if len(hs) > 1 else ""
    _tb(s, 0.8, 0.55, 11.8, 1.0, [cover_title], size=25, color=P_WHITE, bold=True)
    if cover_sub:
        _tb(s, 0.8, 1.45, 11.8, 0.5, [cover_sub], size=13, color=PRGB(0x9F, 0xC5, 0xE0))
    ga = next(((base / v).resolve() for k, v, _ in blocks
               if k == "image" and "graphical_abstract" in v and (base / v).exists()), None)
    if ga: s.shapes.add_picture(str(ga), PIn(1.15), PIn(2.15), width=PIn(11.0))

    # 그림 슬라이드 — 제목·설명을 본문 캡션에서 뽑는다.
    # 파일명을 표에 하드코딩하면 그림 세트가 바뀔 때 조용히 전부 건너뛴다 (실제로 그랬다).
    def _caption_for(idx):
        """이미지 블록 바로 뒤의 문단이 그 그림의 캡션이다."""
        for k, v, _ in blocks[idx + 1: idx + 3]:
            if k == "para" and v.lstrip().startswith("**Figure"):
                return v
        return ""

    def _split_caption(cap):
        """'**Figure 3. 재도킹 대조.** 나머지 설명' → ('재도킹 대조', '나머지 설명')"""
        t = re.sub(r"\*\*", "", cap).strip()
        m = re.match(r"^Figure\s+[\w.]+\.\s*([^.]{2,60})\.\s*(.*)$", t, re.S)
        if not m:
            return (t[:60] or "그림"), ""
        return m.group(1).strip(), re.sub(r"\s+", " ", m.group(2)).strip()[:150]

    for i, (kind, val, _alt) in enumerate(blocks):
        if kind != "image":
            continue
        f = (base / val).resolve()
        if not f.exists():
            continue
        title, sub = _split_caption(_caption_for(i))
        s = new(title, "결과")
        if sub:
            _tb(s, 0.6, 1.32, 12.1, 0.5, [sub], size=10.5, color=P_GREY)
        s.shapes.add_picture(str(f), PIn(0.9), PIn(1.95), width=PIn(11.5))

    # 핵심 결과 표
    tbls = [v for k, v, _ in blocks if k == "table"]
    if len(tbls) >= 2:
        s = new("화합물 물성 실계산값", "결과")
        rows = tbls[1][:7]
        tb = s.shapes.add_table(len(rows), len(rows[0]), PIn(0.7), PIn(1.6),
                                PIn(11.9), PIn(0.4 * len(rows))).table
        for ri, row in enumerate(rows):
            for ci, cell in enumerate(row[:len(rows[0])]):
                c = tb.cell(ri, ci); c.text = clean(cell)
                for pr in c.text_frame.paragraphs:
                    pr.alignment = PP_ALIGN.CENTER
                    for r in pr.runs:
                        r.font.size = PPt(10); r.font.name = KFONT; r.font.bold = (ri == 0)
        _tb(s, 0.7, 1.62 + 0.4 * len(rows) + 0.15, 11.9, 0.6,
            ["QED 내림차순 상위 6건. 값은 전부 RDKit 산출값이며 사람이 적은 수치가 없다."],
            size=11, color=P_GREY)

    # 한계 — 본문이 철회한 것을 덱에서도 밝힌다
    s = new("한계", "이 실행이 말하지 않는 것")
    _tb(s, 0.9, 1.7, 11.6, 3.4,
        ["1.  대조군이 없다.",
         "     문서-전용 조건을 실행하지 않았으므로 '문서보다 코드가 낫다'는 비교 주장은 하지 않는다.",
         "",
         "2.  수정 후에도 selectivity 검사 5개가 이 실행에서는 공허하게 통과한다.",
         "     파이프라인이 SMILES 를 넘기지 않아 입력 의존 검사가 발동하지 않는다.",
         "",
         "3.  단일 실행이며 반복이 없다. 통계 추론을 하지 않았다.",
         "",
         "4.  화합물 10건은 정렬 미지정 조회의 앞부분이고 단일 화학형 계열에 가깝다."],
        size=13, color=P_NAVY)
    _tb(s, 0.9, 5.5, 11.6, 0.9,
        ["게이트 통과는 규약 준수를 뜻할 뿐 화합물이 유망하다는 뜻이 아니다."],
        size=11, color=P_GREY)

    # 결론
    s = new("결론", "정리")
    _tb(s, 0.9, 1.7, 11.6, 3.2,
        ["1.  자세 생성은 작동했고 채점이 작동하지 않았다.",
         "     단, 호출자가 반환값을 검사할 때만이다. 본 실행에서는 4단계 중 3단계만 그랬다.",
         "",
         "2.  강제되었다고 기준이 옳은 것은 아니다.",
         "     이번 실행에서 통과·탈락을 가른 것은 QED 임계 0.5 하나였다.",
         "",
         "3.  임계값은 근거와 함께 명시적으로 정해야 한다.",
         "     스크립트가 스스로 '데모 임계(교육용)'라고 밝힌 값이다."],
        size=14, color=P_NAVY)
    _tb(s, 0.9, 5.4, 11.6, 1.0,
        ["본 보고의 결과 수치는 도구 산출값이다. 서술문과 절 제목은 사람이 작성했다."],
        size=11, color=P_GREY)
    prs.save(out)
    return len(prs.slides._sldIdLst)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--md", default="outputs/report_pde5.md")
    ap.add_argument("--out", default="outputs/docs")
    ap.add_argument("--stem", default="report_pde5", help="산출 파일 이름 줄기")
    a = ap.parse_args()
    md = Path(a.md); out = Path(a.out); out.mkdir(parents=True, exist_ok=True)
    stem = a.stem
    blocks = parse(md.read_text())
    ni, nt = build_docx(blocks, md.parent, out / f"{stem}.docx")
    ns = build_pptx(blocks, md.parent, out / f"{stem}.pptx")
    print(f"  docx: 그림 {ni}장 · 표 {nt}개")
    print(f"  pptx: 슬라이드 {ns}장")

    # 신선도 검사 — 리뷰에서 "절차를 넣었다고 했는데 코드가 없다" 는 지적을 받아 실제로 구현.
    # 본문이나 그림이 문서보다 새로우면 배포본이 철회된 주장을 담을 수 있다.
    # 본문이 실제로 참조하는 그림만 검사한다 — 디렉토리를 하드코딩하면 새 보고서에서 헛돈다.
    figs = [md.parent / src for kind, src, _ in blocks if kind == "image"
            and (md.parent / src).exists()]
    base = max([md.stat().st_mtime] + [f.stat().st_mtime for f in figs])
    stale = [d.name for d in out.glob(f"{stem}*") if d.stat().st_mtime < base]
    if stale:
        raise SystemExit(f"[신선도 실패] 본문·그림보다 오래된 산출물: {stale} — 재생성이 필요합니다.")
    print(f"  신선도: 산출물 {len(list(out.glob(stem + '*')))}건 전부 본문·참조 그림 "
          f"{len(figs)}장보다 최신")


if __name__ == "__main__":
    main()
